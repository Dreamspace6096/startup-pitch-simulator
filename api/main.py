import os
from flask import Flask, request, render_template, jsonify, Response
from werkzeug.utils import secure_filename
import requests
from dotenv import load_dotenv
from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
import tempfile
import traceback
import logging
import json
import queue
import threading

# 导入备选PDF处理方法
try:
    from pdf_processing_alternative import extract_text_from_pdf_without_poppler
    HAS_ALTERNATIVE_PDF = True
except ImportError:
    HAS_ALTERNATIVE_PDF = False

# 设置日志记录
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

# 加载环境变量
load_dotenv()

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = '/tmp/uploads'  # 修改为使用Vercel的/tmp目录
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 限制上传文件大小为50MB

# 创建一个队列来存储状态更新
status_updates = queue.Queue()

# 确保上传目录存在
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

ALLOWED_EXTENSIONS = {'ppt', 'pptx', 'pdf'}

def send_status_update(status):
    """发送状态更新到队列"""
    status_updates.put(json.dumps({"status": status}))

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pptx(file_path):
    """从PPTX文件中提取文本"""
    try:
        prs = Presentation(file_path)
        text_content = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"PPTX处理错误: {str(e)}")
        logger.error(traceback.format_exc())
        raise

def extract_text_from_pdf(file_path):
    """从PDF文件中提取文本"""
    try:
        # 将PDF转换为图像
        images = convert_from_path(file_path)
        text_content = []
        
        # 对每个页面进行OCR
        for image in images:
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            text_content.append(text)
        
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"PDF处理错误: {str(e)}")
        logger.error(traceback.format_exc())
        
        # 如果Poppler方法失败，尝试备选方法
        if HAS_ALTERNATIVE_PDF:
            logger.info("尝试使用备选PDF处理方法...")
            return extract_text_from_pdf_without_poppler(file_path)
        else:
            raise

def generate_questions(text):
    """使用DeepSeek API生成问题"""
    try:
        api_key = os.getenv('DEEPSEEK_API_KEY')
        api_url = os.getenv('DEEPSEEK_API_URL')
        
        if not api_key:
            logger.error("DeepSeek API密钥未配置")
            return {"error": "DeepSeek API密钥未配置"}
        
        if not api_url:
            logger.error("DeepSeek API URL未配置")
            return {"error": "DeepSeek API URL未配置"}
            
        logger.info(f"正在调用DeepSeek API，文本长度：{len(text)}")
        
        prompt = f"""
        作为一位经验丰富的投资人，请根据以下创业项目路演内容，生成5个关键问题。
        每个问题需要简洁明了，直指要害。问题应该覆盖：商业模式、市场规模、竞争优势、团队能力、财务规划等方面。
        请确保问题具有连贯性和深度，能够帮助评估项目的投资价值。
        
        路演内容：
        {text}
        
        请直接列出5个问题，每个问题一行，不要包含序号或其他格式。
        """
        
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }
        
        # 设置请求超时：连接超时为10秒，读取超时为180秒
        try:
            response = requests.post(
                api_url,
                headers=headers,
                json={
                    "model": "deepseek-chat",
                    "messages": [
                        {"role": "system", "content": "你是一位经验丰富的投资人，专门评估创业项目。"},
                        {"role": "user", "content": prompt}
                    ],
                    "temperature": 0.7,
                    "max_tokens": 2000,
                    "top_p": 0.95,
                    "stream": False
                },
                timeout=(10, 180)  # (连接超时, 读取超时)
            )
            
            logger.info(f"API响应状态码：{response.status_code}")
            response_text = response.text
            logger.info(f"API原始响应内容：{response_text[:500]}")
            
            if response.status_code != 200:
                logger.error(f"API请求失败：{response_text}")
                return {"error": f"API请求失败(状态码:{response.status_code}): {response_text[:200]}..."}
            
            try:
                result = response.json()
                if 'choices' in result and len(result['choices']) > 0:
                    questions = result['choices'][0]['message']['content'].strip().split('\n')
                    return {"questions": questions[:5]}
                else:
                    logger.error(f"API响应格式错误：{result}")
                    return {"error": "API响应格式错误，无法提取问题"}
            except json.JSONDecodeError as e:
                logger.error(f"JSON解析错误: {str(e)}, 响应内容: {response_text[:500]}...")
                
                # 尝试手动提取内容（如果响应是纯文本而非JSON）
                if response_text.strip() and len(response_text.strip().split('\n')) >= 3:
                    questions = response_text.strip().split('\n')[:5]
                    logger.info(f"从非JSON响应中提取问题：{questions}")
                    return {"questions": questions}
                    
                return {"error": f"无法解析API响应: {str(e)}，请检查API配置"}
                
        except requests.exceptions.RequestException as e:
            logger.error(f"HTTP请求错误: {str(e)}")
            return {"error": f"HTTP请求错误: {str(e)}"}
            
    except requests.exceptions.ConnectTimeout:
        logger.error("连接API服务器超时")
        return {"error": "连接API服务器超时，请检查网络连接"}
    except requests.exceptions.ReadTimeout:
        logger.error("读取API响应超时")
        return {"error": "读取API响应超时，请稍后重试"}
    except requests.exceptions.ConnectionError as e:
        logger.error(f"连接错误: {str(e)}")
        return {"error": "连接错误，请检查网络连接"}
    except Exception as e:
        logger.error(f"生成问题时发生错误: {str(e)}")
        logger.error(traceback.format_exc())
        return {"error": str(e)}

def analyze_answers(questions, answers):
    """分析回答并生成反馈"""
    try:
        api_key = os.getenv('DEEPSEEK_API_KEY')
        api_url = os.getenv('DEEPSEEK_API_URL')
        
        if not api_key:
            return {"error": "DeepSeek API密钥未配置"}
            
        if not api_url:
            return {"error": "DeepSeek API URL未配置"}
        
        # 将问题和回答组合成对话形式
        qa_pairs = []
        for q_id, question in enumerate(questions):
            answer = answers.get(str(q_id), "")
            qa_pairs.append(f"问：{question}\n答：{answer}")
        qa_text = "\n\n".join(qa_pairs)
        
        prompt = f"""
        作为一位资深投资人，请对以下创业项目路演问答进行专业评估。请从以下几个方面进行分析：

        1. 回答的完整性和逻辑性
        2. 对关键问题的理解深度
        3. 商业模式的可行性论证
        4. 市场认知和竞争分析
        5. 团队能力展示
        6. 财务规划的合理性

        问答内容：
        {qa_text}

        请给出：
        1. 总体评分（满分100分）
        2. 各方面的具体点评
        3. 改进建议
        
        请用简洁的格式输出，使每个点评条目清晰可辨。
        """
        
        # 设置请求超时：连接超时为10秒，读取超时为180秒
        try:
            response = requests.post(
                api_url,
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": "deepseek-chat",
                    "messages": [
                        {"role": "system", "content": "你是一位资深投资人，专门评估创业项目。"},
                        {"role": "user", "content": prompt}
                    ],
                    "temperature": 0.7,
                    "max_tokens": 3000,
                    "top_p": 0.95,
                    "stream": False
                },
                timeout=(10, 180)  # (连接超时, 读取超时)
            )
            
            logger.info(f"分析API响应状态码：{response.status_code}")
            response_text = response.text
            logger.info(f"分析API原始响应内容：{response_text[:500]}")
            
            if response.status_code != 200:
                logger.error(f"分析API请求失败：{response_text}")
                return {"error": f"API请求失败(状态码:{response.status_code}): {response_text[:200]}..."}
                
            try:
                result = response.json()
                if 'choices' in result and len(result['choices']) > 0:
                    feedback = result['choices'][0]['message']['content']
                    return {"feedback": feedback}
                else:
                    logger.error(f"API响应格式错误：{result}")
                    return {"error": "无法生成反馈，API响应格式错误"}
            except json.JSONDecodeError as e:
                logger.error(f"JSON解析错误: {str(e)}, 响应内容: {response_text[:500]}...")
                
                # 如果响应是纯文本而非JSON，直接使用它作为反馈
                if response_text.strip():
                    logger.info("从非JSON响应中提取反馈内容")
                    return {"feedback": response_text}
                    
                return {"error": f"无法解析API响应: {str(e)}，请检查API配置"}
                
        except requests.exceptions.RequestException as e:
            logger.error(f"HTTP请求错误: {str(e)}")
            return {"error": f"HTTP请求错误: {str(e)}"}
            
    except requests.exceptions.ConnectTimeout:
        logger.error("连接API服务器超时")
        return {"error": "连接API服务器超时，请检查网络连接"}
    except requests.exceptions.ReadTimeout:
        logger.error("读取API响应超时")
        return {"error": "读取API响应超时，请稍后重试"}
    except requests.exceptions.ConnectionError as e:
        logger.error(f"连接错误: {str(e)}")
        return {"error": "连接错误，请检查网络连接"}
    except Exception as e:
        logger.error(f"分析回答时发生错误: {str(e)}")
        logger.error(traceback.format_exc())
        return {"error": str(e)}

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload-status')
def status():
    """Server-Sent Events 端点，用于发送状态更新"""
    def generate():
        while True:
            try:
                # 获取状态更新，设置1秒超时
                status_data = status_updates.get(timeout=1)
                yield f"data: {status_data}\n\n"
                if '"status": "completed"' in status_data:
                    break
            except queue.Empty:
                # 发送保持连接的消息
                yield f"data: {json.dumps({'status': 'processing'})}\n\n"
    
    return Response(generate(), mimetype='text/event-stream')

@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({"error": "没有文件被上传"}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "没有选择文件"}), 400
        
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            logger.info(f"保存文件到: {file_path}")
            file.save(file_path)
            
            try:
                # 发送状态：开始读取PDF
                send_status_update("reading_pdf")
                logger.info("开始读取文件")
                
                # 根据文件类型提取文本
                if filename.endswith('.pdf'):
                    logger.info("开始处理PDF文件")
                    text_content = extract_text_from_pdf(file_path)
                else:  # ppt or pptx
                    logger.info("开始处理PPT文件")
                    text_content = extract_text_from_pptx(file_path)
                
                # 发送状态：开始生成问题
                send_status_update("generating_questions")
                logger.info("开始生成问题")
                
                result = generate_questions(text_content)
                
                # 发送状态：完成
                send_status_update("completed")
                logger.info("处理完成")
                
                # 清理临时文件
                os.remove(file_path)
                
                return jsonify(result)
            except Exception as e:
                logger.error(f"处理文件时发生错误: {str(e)}")
                logger.error(traceback.format_exc())
                if os.path.exists(file_path):
                    os.remove(file_path)
                return jsonify({"error": f"处理文件时发生错误: {str(e)}"}), 500
        
        return jsonify({"error": "不支持的文件类型"}), 400
    except Exception as e:
        logger.error(f"上传处理过程中发生错误: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({"error": f"上传处理过程中发生错误: {str(e)}"}), 500

@app.route('/analyze', methods=['POST'])
def analyze():
    try:
        data = request.json
        if not data or 'questions' not in data or 'answers' not in data:
            return jsonify({"error": "缺少必要的数据"}), 400
        
        questions = data['questions']
        answers = data['answers']
        
        if len(questions) != len(answers):
            return jsonify({"error": "问题和答案数量不匹配"}), 400
        
        result = analyze_answers(questions, answers)
        return jsonify(result)
        
    except Exception as e:
        logger.error(f"分析答案时发生错误: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({"error": f"分析答案时发生错误: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', debug=True) 